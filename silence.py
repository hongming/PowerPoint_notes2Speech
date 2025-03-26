import requests
import json
import os
import time
import re
from pptx import Presentation
import argparse
import wave
import struct
import io

# 内建的凭证信息
ALIYUN_APPKEY = "Your key"  # 在此处替换为真实的 AppKey
ALIYUN_TOKEN = "Your token"     # 在此处替换为真实的 Token
SILENCE_DURATION = 2  # 结尾静音时长，单位为秒

class AliyunTTS:
    def __init__(self, appkey=ALIYUN_APPKEY, token=ALIYUN_TOKEN):
        self.appkey = appkey
        self.token = token
        self.url = "https://nls-gateway-cn-shanghai.aliyuncs.com/stream/v1/tts"
        self.headers = {
            'Content-Type': 'application/json',
            'X-NLS-Token': self.token,
        }
    
    def text_to_speech(self, text, output_file, voice='xiaogang', format_type='wav', sample_rate=16000):
        """
        将文本转换为语音
        :param text: 要转换的文本
        :param output_file: 输出文件路径
        :param voice: 发音人，默认为xiaoyun
        :param format_type: 输出格式，默认为wav（注意：改为wav以便易于处理）
        :param sample_rate: 采样率，默认为16000
        :return: 成功返回True，失败返回False
        """
        # 构建请求参数
        data = {
            "appkey": self.appkey,
            "text": text,
            "format": format_type,
            "sample_rate": sample_rate,
            "voice": 'xiaogang',
            "volume": 50,
            "speech_rate": 0,
            "pitch_rate": 0
        }
        
        try:
            # 发送请求
            response = requests.post(self.url, headers=self.headers, data=json.dumps(data))
            
            # 检查响应状态
            if response.status_code == 200:
                # 将二进制内容写入文件
                with open(output_file, 'wb') as f:
                    f.write(response.content)
                print(f"语音合成成功：{output_file}")
                
                # 添加静音
                add_silence_to_wav(output_file, SILENCE_DURATION)
                
                return True
            else:
                print(f"语音合成失败，状态码：{response.status_code}")
                print(f"错误信息：{response.text}")
                return False
        
        except Exception as e:
            print(f"发生错误：{str(e)}")
            return False

def add_silence_to_wav(wav_file, silence_seconds):
    """
    向WAV文件添加指定秒数的静音
    :param wav_file: WAV文件路径
    :param silence_seconds: 静音时长（秒）
    :return: 无
    """
    try:
        # 打开现有WAV文件
        with wave.open(wav_file, 'rb') as wf:
            # 获取音频参数
            channels = wf.getnchannels()
            sample_width = wf.getsampwidth()
            framerate = wf.getframerate()
            n_frames = wf.getnframes()
            
            # 读取所有帧
            frames = wf.readframes(n_frames)
        
        # 计算静音帧数
        silence_frames = int(framerate * silence_seconds)
        
        # 生成静音数据
        silence_data = b'\x00' * channels * sample_width * silence_frames
        
        # 合并原始音频和静音
        new_frames = frames + silence_data
        
        # 创建新的WAV文件
        with wave.open(wav_file, 'wb') as wf:
            wf.setnchannels(channels)
            wf.setsampwidth(sample_width)
            wf.setframerate(framerate)
            wf.writeframes(new_frames)
        
        print(f"已向 {wav_file} 添加 {silence_seconds} 秒静音")
        return True
    
    except Exception as e:
        print(f"添加静音时出错: {str(e)}")
        return False

def extract_notes_from_pptx(pptx_file, output_dir):
    """
    从PowerPoint文件中仅提取每张幻灯片的备注内容并保存为单独的文本文件
    :param pptx_file: PowerPoint文件路径
    :param output_dir: 输出目录
    :return: 无
    """
    try:
        # 创建输出目录
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # 打开PowerPoint文件
        presentation = Presentation(pptx_file)
        
        # 提取的幻灯片计数
        extracted_count = 0
        
        # 遍历每张幻灯片
        for i, slide in enumerate(presentation.slides):
            # 仅提取备注内容
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                
                if notes_text:
                    # 清理文本（去除多余换行符等）
                    notes_text = re.sub(r'\s*\n\s*', ' ', notes_text)  # 将换行替换为空格
                    notes_text = re.sub(r'\s+', ' ', notes_text).strip()  # 清理多余空格
                    
                    # 保存到文本文件，使用三位数编号（001, 002, ...）
                    output_file = os.path.join(output_dir, f"{i+1:03d}.txt")
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(notes_text)
                    
                    print(f"已提取幻灯片 {i+1} 的备注内容")
                    extracted_count += 1
            else:
                print(f"幻灯片 {i+1} 没有备注内容，已跳过")
        
        if extracted_count > 0:
            print(f"成功从 {pptx_file} 中提取了 {extracted_count} 张幻灯片的备注内容")
        else:
            print(f"警告：在 {pptx_file} 中没有找到任何备注内容")
        
        return True
    
    except Exception as e:
        print(f"提取幻灯片备注时出错: {str(e)}")
        return False

def process_text_file(input_file, output_dir):
    """
    处理单个文本文件，将其按照段落分割为多个文本文件
    :param input_file: 输入文本文件
    :param output_dir: 输出目录
    :return: 无
    """
    try:
        # 创建输出目录
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        with open(input_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # 按空行分割段落
        paragraphs = [p.strip() for p in re.split(r'\n\s*\n', content) if p.strip()]
        
        # 保存每个段落为单独的文件
        for i, paragraph in enumerate(paragraphs):
            # 清理文本
            clean_text = re.sub(r'\s*\n\s*', ' ', paragraph)  # 将换行替换为空格
            clean_text = re.sub(r'\s+', ' ', clean_text).strip()  # 清理多余空格
            
            # 保存到文本文件
            output_file = os.path.join(output_dir, f"{i+1:03d}.txt")
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(clean_text)
            
            print(f"已处理第 {i+1} 个段落")
        
        print(f"成功将 {input_file} 分割为 {len(paragraphs)} 个段落")
        return True
    
    except Exception as e:
        print(f"处理文本文件时出错: {str(e)}")
        return False

def generate_audio(text_dir, audio_dir):
    """
    为指定目录中的所有文本文件生成语音
    :param text_dir: 文本文件目录
    :param audio_dir: 输出音频目录
    :return: 无
    """
    try:
        # 创建输出目录
        if not os.path.exists(audio_dir):
            os.makedirs(audio_dir)
        
        # 初始化TTS类（使用内建凭证）
        tts = AliyunTTS()
        
        # 获取所有文本文件
        text_files = [f for f in os.listdir(text_dir) if f.endswith('.txt')]
        # 按数字顺序排序
        text_files.sort(key=lambda x: int(x.split('.')[0]))
        
        if not text_files:
            print(f"警告: 在 {text_dir} 中没有找到任何文本文件")
            return False
        
        for text_file in text_files:
            file_path = os.path.join(text_dir, text_file)
            file_num = text_file.split('.')[0]  # 获取文件编号部分（如"001"）
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read().strip()
            
            if content:  # 确保内容不为空
                # 生成对应的音频文件（WAV格式）
                output_file = os.path.join(audio_dir, f"{file_num}.wav")
                
                # 调用语音合成（带静音）
                success = tts.text_to_speech(content, output_file)
                
                # 避免频繁请求被限制
                if success:
                    time.sleep(1)
        
        return True
    
    except Exception as e:
        print(f"生成音频时出错: {str(e)}")
        return False

def main():
    parser = argparse.ArgumentParser(description='将PowerPoint备注或文本讲解稿转换为语音文件')
    parser.add_argument('--pptx', help='PowerPoint文件路径')
    parser.add_argument('--text', help='文本文件路径')
    parser.add_argument('--output', default='output3', help='输出目录')
    parser.add_argument('--silence', type=int, default=2, help='音频结尾的静音时长(秒)，默认为2秒')
    
    args = parser.parse_args()
    
    if not args.pptx and not args.text:
        print("错误: 请提供PowerPoint文件路径(--pptx)或文本文件路径(--text)")
        return
    
    # 更新静音时长
    global SILENCE_DURATION
    SILENCE_DURATION = args.silence
    
    # 创建所需目录
    text_dir = os.path.join(args.output, "texts")
    audio_dir = os.path.join(args.output, "audios")
    
    if args.pptx:
        print(f"正在从PowerPoint文件 {args.pptx} 提取备注内容...")
        extract_notes_from_pptx(args.pptx, text_dir)
    
    elif args.text:
        print(f"正在处理文本文件 {args.text}...")
        process_text_file(args.text, text_dir)
    
    print(f"正在生成语音文件（每个音频尾部将添加 {args.silence} 秒静音）...")
    generate_audio(text_dir, audio_dir)
    
    print(f"处理完成! 音频文件已保存到: {audio_dir}")

if __name__ == "__main__":
    main()